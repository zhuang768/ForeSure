import os
import copy
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SRC_PATH = "/Users/zhuangzijin/Downloads/ForeSure未然.pptx"
OUT_PATH = "/Users/zhuangzijin/Downloads/ForeSure未然.pptx"
LOCAL_OUT_PATH = "/Users/zhuangzijin/Desktop/FUTUREMODE_x_SITCON_BUILDMODE/ForeSure未然_Completed.pptx"

# Colors
C_DARK_GREEN = RGBColor(0x0E, 0x47, 0x14)    # 0E4714 (Cathay Deep Forest Green)
C_PRIMARY = RGBColor(0x26, 0xA8, 0x62)       # 26A862 (Cathay Green)
C_DARK_TEXT = RGBColor(0x17, 0x21, 0x1C)     # Dark Charcoal Text
C_MUTED = RGBColor(0x4A, 0x55, 0x68)         # Slate Gray Text
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)         # White
C_BG_CARD = RGBColor(0xF9, 0xFA, 0xF9)       # Warm Off-White Card Fill
C_BORDER = RGBColor(0x0E, 0x47, 0x14)        # Deep Green Border
C_BORDER_LIGHT = RGBColor(0xBC, 0xC7, 0xBE)  # Subtle Card Border
C_PM_BLUE = RGBColor(0x1D, 0x4E, 0xD8)       # PM Blue
C_UW_ORANGE = RGBColor(0xC2, 0x41, 0x0C)     # Underwriter Orange
C_AC_GREEN = RGBColor(0x15, 0x80, 0x3D)      # Actuary Green

FONT_HEADING = "Gotham Bold"
FONT_BODY = "Gotham"
FONT_MONO = "Inter"

prs = Presentation(SRC_PATH)
print(f"Initial slide count: {len(prs.slides)}")

# Grab leaf logo element from Slide 2
s2 = prs.slides[1]
leaf_xml = copy.deepcopy(s2.shapes[0].element)

def add_header(slide, tag_text, title_text, subtitle_text=""):
    # Append leaf logo
    slide.shapes._spTree.append(copy.deepcopy(leaf_xml))
    
    # Tag Tracker (e.g. THE SOLUTION)
    tb_tag = slide.shapes.add_textbox(Inches(2.79), Inches(1.05), Inches(10.0), Inches(0.75))
    p = tb_tag.text_frame.paragraphs[0]
    p.text = tag_text
    p.font.name = FONT_HEADING
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = C_DARK_GREEN
    
    # Title & Subtitle
    tb_title = slide.shapes.add_textbox(Inches(1.91), Inches(1.82), Inches(16.5), Inches(1.3))
    tf = tb_title.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = title_text
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = C_DARK_GREEN
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = FONT_BODY
        p2.font.size = Pt(17)
        p2.font.color.rgb = C_MUTED

def create_card(slide, left, top, width, height, border_color=C_BORDER, fill_color=C_WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape

# ==========================================
# SLIDE 3: THE SOLUTION (Tri-Agent Adversarial Pipeline)
# ==========================================
s3 = prs.slides[2]
# Remove old shapes
for s in list(s3.shapes):
    sp = s._element
    sp.getparent().remove(sp)

add_header(s3, "THE SOLUTION", "Tri-Agent Adversarial Consensus Pipeline", 
           "Moving from raw telemetry to grounded policy through continuous dialectic checks")

# 3 Horizontal Process Cards
# Card 1: Ingestion & RAG
create_card(s3, Inches(0.6), Inches(3.6), Inches(5.6), Inches(7.0))
tb1 = s3.shapes.add_textbox(Inches(0.9), Inches(3.9), Inches(5.0), Inches(6.4))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "STEP 01 // TELEMETRY"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_PM_BLUE

p = tf1.add_paragraph()
p.text = "News Radar & ChromaDB"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items1 = [
    "· Live Google News RSS automatically captures emerging peril events (climate, cyber, grid outages)",
    "· ChromaDB dense semantic vector search indexes 30 Cathay Century Insurance policies",
    "· Cosine similarity instantly pinpoints unserved policy gaps (< 3 seconds)",
    "· Contextual gap analysis is compiled into structured prompts for the PM Agent"
]
for item in items1:
    p = tf1.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(10)

# Card 2: Adversarial Arena
create_card(s3, Inches(6.6), Inches(3.6), Inches(6.8), Inches(7.0))
tb2 = s3.shapes.add_textbox(Inches(6.9), Inches(3.9), Inches(6.2), Inches(6.4))
tf2 = tb2.text_frame
tf2.word_wrap = True

p = tf2.paragraphs[0]
p.text = "STEP 02 // DIALECTIC ARENA"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_UW_ORANGE

p = tf2.add_paragraph()
p.text = "PM vs Underwriter Debate"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

p = tf2.add_paragraph()
p.text = "Product Manager (Expansion Bias)"
p.font.name = FONT_HEADING
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = C_PM_BLUE
p.space_before = Pt(8)

p = tf2.add_paragraph()
p.text = "· Drafts high-appeal parametric structures to capture market gaps"
p.font.name = FONT_BODY
p.font.size = Pt(16)
p.font.color.rgb = C_DARK_TEXT

p = tf2.add_paragraph()
p.text = "▲▼ Multi-Round Adversarial Cross-Examination"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN
p.space_before = Pt(12)

p = tf2.add_paragraph()
p.text = "Senior Underwriter (Risk Shield)"
p.font.name = FONT_HEADING
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = C_UW_ORANGE
p.space_before = Pt(8)

p = tf2.add_paragraph()
p.text = "· Ruthlessly challenges moral hazards, fraud loopholes & catastrophe aggregation"
p.font.name = FONT_BODY
p.font.size = Pt(16)
p.font.color.rgb = C_DARK_TEXT

p = tf2.add_paragraph()
p.text = "· Enforces strict exclusion clauses and verifies objective parametric data sources"
p.font.name = FONT_BODY
p.font.size = Pt(16)
p.font.color.rgb = C_DARK_TEXT

# Card 3: Actuarial Convergence
create_card(s3, Inches(13.8), Inches(3.6), Inches(5.6), Inches(7.0))
tb3 = s3.shapes.add_textbox(Inches(14.1), Inches(3.9), Inches(5.0), Inches(6.4))
tf3 = tb3.text_frame
tf3.word_wrap = True

p = tf3.paragraphs[0]
p.text = "STEP 03 // CONVERGENCE"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_AC_GREEN

p = tf3.add_paragraph()
p.text = "Actuarial Structuring"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items3 = [
    "· Reconciles commercial growth with strict underwriting solvency boundaries",
    "· Solvency II / TW-ICS 99.5% capital margin with 1.2x - 3.0x dynamic safety markup",
    "· Pure algorithmic Function Calling outputs 12 structured Chinese & English fields",
    "· Generates 32-byte SHA-256 fingerprint for Ethereum Sepolia smart contract"
]
for item in items3:
    p = tf3.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(10)

print("Slide 3 built successfully.")

# ==========================================
# SLIDE 4: ACTUARIAL RIGOR (Empirical Data Modeling)
# ==========================================
s4 = prs.slides[3]
for s in list(s4.shapes):
    sp = s._element
    sp.getparent().remove(sp)

add_header(s4, "ACTUARIAL RIGOR", "67-Year Empirical NFA Disaster Loss Modeling",
           "Grounded in real government disaster statistics, Poisson event frequencies, and Solvency II margins")

# Two large asymmetric panels
# Left Panel: Empirical Historical Foundation
create_card(s4, Inches(0.8), Inches(3.6), Inches(9.0), Inches(7.0))
tb4_l = s4.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(8.2), Inches(6.4))
tf4_l = tb4_l.text_frame
tf4_l.word_wrap = True

p = tf4_l.paragraphs[0]
p.text = "EMPIRICAL HISTORICAL BASELINE"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

p = tf4_l.add_paragraph()
p.text = "67 Years of Disaster Data (1958 - 2025)"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items4_l = [
    "· National Fire Agency (NFA) Historical Records: 67 years of Taiwan natural disaster statistics (typhoons, floods, earthquakes).",
    "· Post-1995 Severe Event Filtering: Evaluates events with >= 50 destroyed households post-1995 building code overhaul to reflect current exposure.",
    "· Poisson Annual Frequency Modeling: Calculates empirical annual arrival probability λ without artificial curve-smoothing bias.",
    "· Taiwan Earthquake Insurance Benchmark: NT$ 1,500,000 full-loss benefit baseline from Residential Earthquake Basic Insurance.",
    "· Preserving True Tail Risks: Captures severe historical events (e.g. 921 Earthquake) to test heavy-tail exposure accurately."
]
for item in items4_l:
    p = tf4_l.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(8)

# Right Panel: Transparent Attribution
create_card(s4, Inches(10.2), Inches(3.6), Inches(9.0), Inches(7.0))
tb4_r = s4.shapes.add_textbox(Inches(10.6), Inches(3.9), Inches(8.2), Inches(6.4))
tf4_r = tb4_r.text_frame
tf4_r.word_wrap = True

p = tf4_r.paragraphs[0]
p.text = "TRANSPARENT EVIDENCE ATTRIBUTION"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_PRIMARY

p = tf4_r.add_paragraph()
p.text = "Real Statistics vs Explicit Assumptions"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items4_r = [
    "· Real Statistics Badge (真實統計): Typhoon, flood, and earthquake probabilities are tied directly to official government archives.",
    "· Explicit Assumption Disclosure (假設值揭露): Emerging perils without official loss history (ransomware, cloud outages) are strictly flagged as assumptions.",
    "· Dynamic Safety Markup: 1.2x to 3.0x markup scaling with event frequency to absorb parameter uncertainty and operational expenses.",
    "· Solvency II / TW-ICS Capital Margins: Dynamic 99.5% capital adequacy stress testing ensures financial solvency.",
    "· Complete Audit Transparency: Actuarial basis metadata attached to every single figure, preventing misleading quotes."
]
for item in items4_r:
    p = tf4_r.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(8)

print("Slide 4 built successfully.")

# ==========================================
# SLIDE 5: TRUST & AUDIT (Grounding Check & Blockchain)
# ==========================================
s5 = prs.slides[4]
for s in list(s5.shapes):
    sp = s._element
    sp.getparent().remove(sp)

add_header(s5, "TRUST & AUDIT", "Deterministic Grounding & Ethereum Sepolia Attestation",
           "Eliminating hallucinations with non-LLM algorithmic rules and immutable smart contract proofs")

# Left Panel: Non-LLM Grounding Checker
create_card(s5, Inches(0.8), Inches(3.6), Inches(9.0), Inches(7.0))
tb5_l = s5.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(8.2), Inches(6.4))
tf5_l = tb5_l.text_frame
tf5_l.word_wrap = True

p = tf5_l.paragraphs[0]
p.text = "DEFENSE LAYER 01 // GROUNDING CHECK"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

p = tf5_l.add_paragraph()
p.text = "Pure Algorithmic Anti-Hallucination (v1.1)"
p.font.name = FONT_HEADING
p.font.size = Pt(23)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items5_l = [
    "· 100% Deterministic & Non-LLM: Does not rely on model self-evaluation. 100 tests produce 100 identical audit verdicts.",
    "· Unsupported Number Audit (High Severity): Every number in market gap and business logic must trace to the actuarial engine, news text, or policy database (2% tolerance).",
    "· Unverified Citation Trap (High Severity): Any 'according to X' claim must exist verbatim in ingested evidence.",
    "· Missing Disclosure Flag (Medium Severity): Actuarial assumptions must explicitly contain 'estimated' or 'assumed'.",
    "· Tamper-Proof Audit Status: Grounding verdict (pass / warn / fail) and flag counts are sealed into the decision hash."
]
for item in items5_l:
    p = tf5_l.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(8)

# Right Panel: Ethereum Sepolia Notary
create_card(s5, Inches(10.2), Inches(3.6), Inches(9.0), Inches(7.0))
tb5_r = s5.shapes.add_textbox(Inches(10.6), Inches(3.9), Inches(8.2), Inches(6.4))
tf5_r = tb5_r.text_frame
tf5_r.word_wrap = True

p = tf5_r.paragraphs[0]
p.text = "DEFENSE LAYER 02 // BLOCKCHAIN PROOF"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_PRIMARY

p = tf5_r.add_paragraph()
p.text = "Ethereum Sepolia Smart Contract Notary"
p.font.name = FONT_HEADING
p.font.size = Pt(23)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items5_r = [
    "· 32-Byte SHA-256 Decision Fingerprint: 13 canonical fields (decision ID, news trigger, pricing, grounding verdict) hashed into AuditRegistry.sol.",
    "· Zero Business Secret Leakage: Only the cryptographic hash is published on-chain; proprietary trade secrets remain in internal storage.",
    "· Append-Only Immutable Registry: The smart contract contains no update or delete functions, ensuring regulatory auditability.",
    "· Live Destructive Tamper Test: Judges can alter any single digit on the live dashboard—on-chain verification immediately fails with red alert.",
    "· Public Verification on Etherscan: Directly verifiable by regulators and auditors on the public Ethereum Sepolia testnet."
]
for item in items5_r:
    p = tf5_r.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(8)

print("Slide 5 built successfully.")

# ==========================================
# SLIDE 6: ENTERPRISE DEPLOYMENT (Architecture & Live Ecosystem)
# ==========================================
s6 = prs.slides[5]
for s in list(s6.shapes):
    sp = s._element
    sp.getparent().remove(sp)

add_header(s6, "ENTERPRISE DEPLOYMENT", "Apigee Gateway, Cloudflare Edge & Bilingual Dossier",
           "From rapid hackathon prototype to production-grade financial enterprise architecture")

# 3 Architecture Columns
create_card(s6, Inches(0.6), Inches(3.6), Inches(5.6), Inches(7.0))
tb6_1 = s6.shapes.add_textbox(Inches(0.9), Inches(3.9), Inches(5.0), Inches(6.4))
tf6_1 = tb6_1.text_frame
tf6_1.word_wrap = True

p = tf6_1.paragraphs[0]
p.text = "TIER 01 // SECURITY GATEWAY"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

p = tf6_1.add_paragraph()
p.text = "Cathay Apigee Ready"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items6_1 = [
    "· Enterprise JWT Bearer Authentication (HTTPBearer) protecting pipeline trigger endpoints",
    "· Client IP Rate Limiting (30 requests/minute) guarding against DDoS and quota exhaustion",
    "· 12-Stage Server-Sent Events (SSE) live pipeline stream at /api/v1/runs/{id}/events",
    "· Full compliance with Cathay Financial Holdings API security specifications"
]
for item in items6_1:
    p = tf6_1.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(10)

create_card(s6, Inches(6.6), Inches(3.6), Inches(6.8), Inches(7.0))
tb6_2 = s6.shapes.add_textbox(Inches(6.9), Inches(3.9), Inches(6.2), Inches(6.4))
tf6_2 = tb6_2.text_frame
tf6_2.word_wrap = True

p = tf6_2.paragraphs[0]
p.text = "TIER 02 // GLOBAL EDGE UI"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_PRIMARY

p = tf6_2.add_paragraph()
p.text = "Cloudflare Edge Production"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items6_2 = [
    "· Next.js 16 + React 19 + TypeScript deployed globally on Cloudflare Pages edge CDN",
    "· 0-ms Instant Loading Archive: 20 historical proposals preloaded for offline judge browsing",
    "· Financial Instrument Aesthetic: Zero emojis, dark/light OLED theme, and bilingual ZH/EN toggle",
    "· Live decision replay and interactive on-chain tamper detection console"
]
for item in items6_2:
    p = tf6_2.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(10)

create_card(s6, Inches(13.8), Inches(3.6), Inches(5.6), Inches(7.0))
tb6_3 = s6.shapes.add_textbox(Inches(14.1), Inches(3.9), Inches(5.0), Inches(6.4))
tf6_3 = tb6_3.text_frame
tf6_3.word_wrap = True

p = tf6_3.paragraphs[0]
p.text = "TIER 03 // FILING DOSSIER"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_AC_GREEN

p = tf6_3.add_paragraph()
p.text = "Automated Bilingual Word"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items6_3 = [
    "· 1-Click Automated .docx Dossier Generation for internal insurance committee review",
    "· Side-by-Side Traditional Chinese and English text across all six policy clauses",
    "· Clickable Hyperlinks directly connecting back to original trigger news sources",
    "· Embedded Grounding Verdict badge, Etherscan transaction link, and hash seal"
]
for item in items6_3:
    p = tf6_3.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(10)

print("Slide 6 built successfully.")

# ==========================================
# SLIDE 7: LIVE DEMO & IMPACT (Strategic Value & Team)
# ==========================================
s7 = prs.slides[6]
for s in list(s7.shapes):
    sp = s._element
    sp.getparent().remove(sp)

add_header(s7, "LIVE DEMO & IMPACT", "The Actuary's AI Decision Co-Pilot in Production",
           "AI accelerates the draft in 85s. Human experts sign off. The blockchain guarantees permanent trust.")

# Split Screen: Left Demo Guide, Right Team & Mission
create_card(s7, Inches(0.8), Inches(3.6), Inches(9.2), Inches(7.0))
tb7_l = s7.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(8.4), Inches(6.4))
tf7_l = tb7_l.text_frame
tf7_l.word_wrap = True

p = tf7_l.paragraphs[0]
p.text = "LIVE PRODUCTION ACCESS"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

p = tf7_l.add_paragraph()
p.text = "Interactive Demo Playbook"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items7_l = [
    "· Production URL: https://atlas-insurance-dashboard.pages.dev/",
    "· Step 1: System Intro (/) — Explore Tri-Agent architecture & real-time telemetry strip.",
    "· Step 2: Live Analysis (/generator) — 85-second live news crawl, 3-way debate, and on-chain sealing.",
    "· Step 3: Historical Vault (/history) — Browse 20 pre-audited records with instant search and debate replay.",
    "· Step 4: Audit & Tamper Test — Click 'Verify' for contract match; click 'Tamper Test' to alter probability and watch the blockchain fail instantly."
]
for item in items7_l:
    p = tf7_l.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(8)

create_card(s7, Inches(10.4), Inches(3.6), Inches(8.8), Inches(7.0))
tb7_r = s7.shapes.add_textbox(Inches(10.8), Inches(3.9), Inches(8.0), Inches(6.4))
tf7_r = tb7_r.text_frame
tf7_r.word_wrap = True

p = tf7_r.paragraphs[0]
p.text = "STRATEGIC MISSION & TEAM"
p.font.name = FONT_HEADING
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_PRIMARY

p = tf7_r.add_paragraph()
p.text = "Team ForeSure 未然"
p.font.name = FONT_HEADING
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = C_DARK_GREEN

items7_r = [
    "· Team Members: TZU-CHIN CHUANG (莊子進) · WEN-HAN LEE (李文瀚)",
    "· Event: FUTUREMODE x SITCON BUILDMODE Hackathon 2026",
    "· Core Philosophy: We do NOT replace actuaries. We empower them to start from a rigorous, data-backed draft rather than a blank sheet.",
    "· Regulatory Alignment: Strictly complies with Taiwan Insurance Bureau filing norms; human appointed actuaries retain final sign-off authority.",
    "· Blockchain's Real Role: Not for volatile token payments, but as an immutable public notary ensuring AI transparency and accountability."
]
for item in items7_r:
    p = tf7_r.add_paragraph()
    p.text = item
    p.font.name = FONT_BODY
    p.font.size = Pt(17)
    p.font.color.rgb = C_DARK_TEXT
    p.space_before = Pt(8)

print("Slide 7 built successfully.")

# ==========================================
# TRIM UNUSED TEMPLATE SLIDES (Slides 8 to 16)
# ==========================================
# Keep exactly 7 high-impact slides
slide_id_list = prs.slides._sldIdLst
while len(slide_id_list) > 7:
    slide_to_delete = slide_id_list[-1]
    rId = slide_to_delete.rId
    prs.part.drop_rel(rId)
    slide_id_list.remove(slide_to_delete)

print(f"Final slide count: {len(prs.slides)}")

# Save to Downloads and local workspace
prs.save(OUT_PATH)
prs.save(LOCAL_OUT_PATH)
print(f"Saved successfully to:\n  - {OUT_PATH}\n  - {LOCAL_OUT_PATH}")

